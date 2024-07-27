from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Project Title, Guide and Student details
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
title.text = "Deep Learning of Facial Depth Maps for Obstructive Sleep Apnea Prediction"
subtitle = slide_1.placeholders[1]
subtitle.text = "Guide: [Guide Name]\nStudent: [Student Name]"

# Slide 2: Abstract
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_2.shapes.title
title.text = "Abstract"
content = slide_2.placeholders[1]
content.text = (
    "Obstructive Sleep Apnea (OSA) occurs when the airway is repeatedly obstructed during sleep due to relaxation of the tongue and airway muscles. "
    "Indicators of OSA include snoring, poor sleep quality, and waking up unrefreshed. OSA diagnosis is costly, leading to many undiagnosed cases. "
    "Previous research links facial morphology with OSA. This project investigates using deep learning on facial depth maps for OSA diagnosis, "
    "achieving around 69% validation accuracy using transfer learning."
)

# Slide 3: Introduction – Aim, Objective, Scope
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_3.shapes.title
title.text = "Introduction – Aim, Objective, Scope"
content = slide_3.placeholders[1]
content.text = (
    "Aim: To develop a deep learning model for predicting OSA using facial depth maps.\n\n"
    "Objective:\n"
    "- To apply deep learning techniques on facial depth maps.\n"
    "- To improve the accuracy of OSA prediction.\n"
    "- To compare the new method with existing diagnostic techniques.\n\n"
    "Scope: This study focuses on the correlation between facial morphology and OSA, using depth maps for more detailed analysis compared to 2D images."
)

# Slide 4: Existing System and its limitations
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_4.shapes.title
title.text = "Existing System and its limitations"
content = slide_4.placeholders[1]
content.text = (
    "Existing System:\n"
    "- Uses imaging techniques like MRI and dental X-rays to assess upper airway structures.\n"
    "- Employs volumetric analysis and tomography for craniofacial measurements.\n\n"
    "Limitations:\n"
    "- Low accuracy and efficiency.\n"
    "- Manual landmark detection is time-consuming and dependent on the operator's expertise."
)

# Slide 5: Proposed System and its advantages
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_5.shapes.title
title.text = "Proposed System and its advantages"
content = slide_5.placeholders[1]
content.text = (
    "Proposed System:\n"
    "- Utilizes pre-trained CNN models (e.g., VGGFace, Pose-Aware CNN Models) for transfer learning with facial depth maps.\n"
    "- Fine-tunes networks for end-to-end classification of OSA from facial depth maps.\n\n"
    "Advantages:\n"
    "- Higher accuracy and efficiency.\n"
    "- Automates landmark detection, reducing time and dependency on manual expertise."
)

# Slide 6: Software and Hardware requirements
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_6.shapes.title
title.text = "Software and Hardware requirements"
content = slide_6.placeholders[1]
content.text = (
    "Hardware Requirements:\n"
    "- System: i3 or above.\n"
    "- RAM: 4 GB.\n"
    "- Hard Disk: 40 GB.\n\n"
    "Software Requirements:\n"
    "- Operating System: Windows 8 or above.\n"
    "- Coding Language: Python\n"
    "- Image Format: JPEG\n"
    "- Dimensions: 221 x 228"
)

# Slide 7: Software Architecture (Block / ER/ DFD/ UML Diagrams)
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_7.shapes.title
title.text = "Software Architecture"
content = slide_7.placeholders[1]
content.text = "Block Diagram of CNN models (VGG19, AlexNet)\nUML Diagrams including Use Case, Class, Sequence, and Collaboration diagrams."

# Slide 8: Modules
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_8.shapes.title
title.text = "Modules"
content = slide_8.placeholders[1]
content.text = (
    "Modules:\n"
    "- Upload OSA face dataset\n"
    "- Preprocess dataset\n"
    "- Build VGG19 model\n"
    "- Upload test data & predict OSA\n"
    "- Accuracy comparison graph"
)

# Save the presentation
prs.save("Deep_Learning_OSA_Prediction_Presentation.pptx")
